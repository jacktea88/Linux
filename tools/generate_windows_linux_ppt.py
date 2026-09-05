from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.util import Inches, Pt


def set_bg(slide, rgb):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*rgb)


def add_textbox(slide, x, y, w, h, text, size=18, bold=False, color=(17, 24, 39), align_center=True):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "Noto Sans TC"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)
    if align_center:
        p.alignment = 1
    return box


def add_box(slide, x, y, w, h, text, fill_rgb=(248, 250, 252), line_rgb=(51, 65, 85), text_size=16):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*fill_rgb)
    shape.line.color.rgb = RGBColor(*line_rgb)
    shape.line.width = Pt(1.8)

    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "Noto Sans TC"
    run.font.size = Pt(text_size)
    run.font.color.rgb = RGBColor(17, 24, 39)
    p.alignment = 1
    return shape


def add_arrow(slide, x, y, w, h, color=(6, 182, 212)):
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = RGBColor(*color)
    arrow.line.color.rgb = RGBColor(14, 116, 144)
    arrow.line.width = Pt(1.2)
    return arrow


def add_line(slide, x1, y1, x2, y2, color=(6, 182, 212), width=1.5, dashed=False):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    line.line.color.rgb = RGBColor(*color)
    line.line.width = Pt(width)
    if dashed:
        line.line.dash_style = 2
    return line


def build_ppt(output_path: Path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, (248, 250, 252))

    add_textbox(
        slide, 0.4, 0.15, 12.6, 0.55,
        "Windows 與 Linux 的系統架構比較",
        size=34,
        bold=True,
        color=(15, 23, 42),
    )
    add_textbox(
        slide, 0.4, 0.65, 12.6, 0.35,
        "從硬體資源到核心與應用層",
        size=18,
        color=(71, 85, 105),
    )

    add_textbox(slide, 0.35, 1.1, 2.8, 0.35, "電腦硬體", size=20, bold=True)
    add_textbox(slide, 3.45, 1.1, 3.6, 0.35, "作業系統", size=20, bold=True)
    add_textbox(slide, 7.35, 1.1, 5.5, 0.35, "系統分層", size=20, bold=True)

    hw_cpu = add_box(slide, 0.55, 1.7, 2.4, 0.9, "CPU")
    hw_ram = add_box(slide, 0.55, 2.8, 2.4, 0.9, "RAM")
    hw_hdd = add_box(slide, 0.55, 3.9, 2.4, 0.9, "HDD")

    win_box = add_box(
        slide, 3.55, 1.7, 3.2, 1.4,
        "Windows\nC: / D:",
        fill_rgb=(236, 254, 255),
        line_rgb=(6, 182, 212),
    )
    linux_box = add_box(
        slide, 3.55, 3.35, 3.2, 1.8,
        "Linux\n/  /home  /etc  /var  ...",
        fill_rgb=(240, 253, 244),
        line_rgb=(22, 163, 74),
    )

    stack_top = add_box(slide, 7.55, 1.7, 5.2, 1.2, "Applications\nServer Apps / X-Window Apps")
    stack_mid = add_box(slide, 7.55, 3.0, 5.2, 0.8, "OS Utilities / Applications")
    stack_kernel = add_box(slide, 7.55, 3.9, 5.2, 0.8, "Kernel")
    stack_hw = add_box(slide, 7.55, 4.8, 5.2, 0.8, "Hardware")

    add_arrow(slide, 6.95, 2.1, 0.5, 0.35)
    add_arrow(slide, 6.95, 3.55, 0.5, 0.35)

    add_line(slide, 2.95, 2.15, 3.55, 2.15)
    add_line(slide, 2.95, 3.25, 3.55, 2.35)
    add_line(slide, 2.95, 4.35, 3.55, 4.2)

    add_line(slide, 2.95, 2.15, 3.55, 4.2)
    add_line(slide, 2.95, 3.25, 3.55, 4.2)
    add_line(slide, 2.95, 4.35, 3.55, 4.2)

    add_textbox(
        slide, 3.8, 5.35, 3.1, 0.35,
        "符合 FHS 標準規範",
        size=14,
        color=(21, 128, 61),
    )
    add_line(slide, 5.15, 5.2, 4.95, 5.0, color=(22, 163, 74), width=1.3, dashed=True)

    note = add_textbox(
        slide, 0.5, 6.5, 12.3, 0.45,
        "重點：Linux 使用單一目錄樹管理資源，並常遵循 FHS，方便維護與部署。",
        size=14,
        color=(30, 41, 59),
        align_center=False,
    )
    note.text_frame.margin_left = Inches(0.05)

    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide2, (248, 250, 252))

    add_textbox(
        slide2, 0.4, 0.2, 12.6, 0.55,
        "Linux FHS 目錄結構詳解",
        size=34,
        bold=True,
        color=(15, 23, 42),
    )
    add_textbox(
        slide2, 0.4, 0.72, 12.6, 0.32,
        "Filesystem Hierarchy Standard：讓系統結構一致、可維護、可移植",
        size=16,
        color=(71, 85, 105),
    )

    add_box(
        slide2, 0.55, 1.25, 3.0, 4.9,
        "/\n|-- /bin\n|-- /etc\n|-- /home\n|-- /usr\n|-- /var\n`-- /tmp",
        fill_rgb=(240, 253, 244),
        line_rgb=(22, 163, 74),
        text_size=18,
    )

    add_box(
        slide2, 3.95, 1.25, 8.8, 0.95,
        "/bin：基本使用者命令（如 ls、cp、cat）",
        fill_rgb=(248, 250, 252),
        line_rgb=(51, 65, 85),
        text_size=16,
    )
    add_box(
        slide2, 3.95, 2.3, 8.8, 0.95,
        "/etc：系統與服務設定檔（如 sshd_config）",
        fill_rgb=(248, 250, 252),
        line_rgb=(51, 65, 85),
        text_size=16,
    )
    add_box(
        slide2, 3.95, 3.35, 8.8, 0.95,
        "/home：一般使用者家目錄與個人檔案",
        fill_rgb=(248, 250, 252),
        line_rgb=(51, 65, 85),
        text_size=16,
    )
    add_box(
        slide2, 3.95, 4.4, 8.8, 0.95,
        "/usr：應用程式與共享資源（多數套件安裝位置）",
        fill_rgb=(248, 250, 252),
        line_rgb=(51, 65, 85),
        text_size=16,
    )
    add_box(
        slide2, 3.95, 5.45, 8.8, 0.95,
        "/var：變動資料（log、cache、spool）；/tmp：暫存檔",
        fill_rgb=(248, 250, 252),
        line_rgb=(51, 65, 85),
        text_size=16,
    )

    add_textbox(
        slide2, 0.55, 6.55, 12.2, 0.35,
        "教學重點：找設定先看 /etc，找日誌先看 /var/log，找使用者檔案看 /home。",
        size=14,
        color=(30, 41, 59),
        align_center=False,
    )

    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide3, (248, 250, 252))

    add_textbox(
        slide3, 0.4, 0.2, 12.6, 0.55,
        "授課節奏與動畫順序（90 秒）",
        size=32,
        bold=True,
        color=(15, 23, 42),
    )
    add_textbox(
        slide3, 0.4, 0.72, 12.6, 0.32,
        "建議先講架構，再帶 FHS，最後用三句口訣收斂",
        size=16,
        color=(71, 85, 105),
    )

    add_box(
        slide3, 0.55, 1.25, 6.1, 4.55,
        "動畫順序\n\n"
        "1) 標題與副標淡入（0-8s）\n"
        "2) 硬體區塊依序出現（8-20s）\n"
        "3) Windows / Linux 方塊出現（20-35s）\n"
        "4) 流向箭頭補齊（35-55s）\n"
        "5) 右側分層由下往上（55-75s）\n"
        "6) FHS 註解出現收尾（75-90s）",
        fill_rgb=(236, 254, 255),
        line_rgb=(6, 182, 212),
        text_size=16,
    )

    add_box(
        slide3, 6.95, 1.25, 5.8, 4.55,
        "90 秒口白提要\n\n"
        "同一台硬體可由不同作業系統管理。\n"
        "Windows 常用磁碟代號；Linux 使用單一目錄樹。\n"
        "Kernel 位於中間，負責軟硬體溝通。\n"
        "Linux 常遵循 FHS，因此部署與維護更一致。\n"
        "口訣：找設定看 /etc，找日誌看 /var/log，\n"
        "找使用者檔案看 /home。",
        fill_rgb=(240, 253, 244),
        line_rgb=(22, 163, 74),
        text_size=16,
    )

    add_textbox(
        slide3, 0.55, 6.05, 12.2, 0.8,
        "授課技巧：每個步驟只停留一個重點，搭配滑鼠圈選關鍵詞，學生理解率會更高。",
        size=14,
        color=(30, 41, 59),
        align_center=False,
    )

    prs.save(output_path)


def build_student_ppt(output_path: Path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, (248, 250, 252))

    add_textbox(
        slide, 0.4, 0.15, 12.6, 0.65,
        "Windows vs Linux 系統架構",
        size=42,
        bold=True,
        color=(15, 23, 42),
    )
    add_textbox(
        slide, 0.4, 0.78, 12.6, 0.35,
        "先看硬體，再看作業系統，最後看應用層",
        size=22,
        color=(71, 85, 105),
    )

    add_textbox(slide, 0.35, 1.2, 2.8, 0.45, "硬體", size=28, bold=True)
    add_textbox(slide, 3.45, 1.2, 3.6, 0.45, "作業系統", size=28, bold=True)
    add_textbox(slide, 7.35, 1.2, 5.5, 0.45, "系統分層", size=28, bold=True)

    add_box(slide, 0.55, 1.95, 2.4, 1.0, "CPU", text_size=24)
    add_box(slide, 0.55, 3.1, 2.4, 1.0, "RAM", text_size=24)
    add_box(slide, 0.55, 4.25, 2.4, 1.0, "HDD", text_size=24)

    add_box(
        slide, 3.55, 1.95, 3.2, 1.6,
        "Windows\nC: / D:",
        fill_rgb=(236, 254, 255),
        line_rgb=(6, 182, 212),
        text_size=24,
    )
    add_box(
        slide, 3.55, 3.75, 3.2, 1.8,
        "Linux\n/ /home /etc /var",
        fill_rgb=(240, 253, 244),
        line_rgb=(22, 163, 74),
        text_size=22,
    )

    add_box(slide, 7.55, 1.95, 5.2, 1.4, "Applications", text_size=26)
    add_box(slide, 7.55, 3.45, 5.2, 0.95, "OS Utilities", text_size=24)
    add_box(slide, 7.55, 4.5, 5.2, 0.95, "Kernel", text_size=24)
    add_box(slide, 7.55, 5.55, 5.2, 0.95, "Hardware", text_size=24)

    add_arrow(slide, 6.95, 2.3, 0.5, 0.42)
    add_arrow(slide, 6.95, 4.2, 0.5, 0.42)

    add_textbox(
        slide, 0.55, 6.75, 12.2, 0.35,
        "重點：Linux 使用單一目錄樹，常遵循 FHS。",
        size=20,
        color=(21, 128, 61),
        align_center=False,
    )

    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide2, (248, 250, 252))

    add_textbox(
        slide2, 0.4, 0.2, 12.6, 0.65,
        "Linux FHS 快速記憶",
        size=42,
        bold=True,
        color=(15, 23, 42),
    )

    add_box(
        slide2, 0.8, 1.2, 12.0, 1.0,
        "/etc = 設定檔    /var/log = 日誌    /home = 使用者檔案",
        fill_rgb=(240, 253, 244),
        line_rgb=(22, 163, 74),
        text_size=27,
    )
    add_box(
        slide2, 0.8, 2.45, 12.0, 1.0,
        "/bin = 基本命令    /usr = 多數程式與資源",
        fill_rgb=(236, 254, 255),
        line_rgb=(6, 182, 212),
        text_size=27,
    )
    add_box(
        slide2, 0.8, 3.7, 12.0, 1.0,
        "口訣：找設定看 /etc，找日誌看 /var/log，找個人檔案看 /home",
        fill_rgb=(248, 250, 252),
        line_rgb=(51, 65, 85),
        text_size=26,
    )

    add_textbox(
        slide2, 0.8, 5.1, 12.0, 1.7,
        "練習題：\n1. SSH 設定檔通常在哪個目錄？\n2. 系統錯誤日誌通常去哪裡找？\n3. 使用者下載檔案通常放在哪裡？",
        size=24,
        color=(30, 41, 59),
        align_center=False,
    )

    prs.save(output_path)


if __name__ == "__main__":
    root = Path(__file__).resolve().parents[1]
    out_teacher = root / "PPT" / "Windows_Linux_教材版.pptx"
    out_student = root / "PPT" / "Windows_Linux_學生版.pptx"
    try:
        build_ppt(out_teacher)
        print(f"Generated: {out_teacher}")
    except PermissionError:
        print(f"Skipped (file is open): {out_teacher}")
    build_student_ppt(out_student)
    print(f"Generated: {out_student}")
